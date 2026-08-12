from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import requests
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from dotenv import load_dotenv
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import google.generativeai as genai
import json

load_dotenv()

HF_API_TOKEN = os.environ.get("HF_API_TOKEN")
HF_SUMMARY_URL = "https://router.huggingface.co/hf-inference/models/facebook/bart-large-cnn"

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)
gemini_model = genai.GenerativeModel('gemini-flash-latest')

PIXABAY_API_KEY = os.environ.get("PIXABAY_API_KEY")

LANGUAGETOOL_URL = "https://api.languagetool.org/v2/check"

app = Flask(__name__)
CORS(app)

similarity_model = SentenceTransformer('all-MiniLM-L6-v2')


@app.route("/check-grammar", methods=["POST"])
def check_grammar():
    data = request.get_json()
    text = data.get("text", "")
    language = data.get("language", "en-US")

    if not text.strip():
        return jsonify({"error": "Text is empty"}), 400

    try:
        response = requests.post(
            LANGUAGETOOL_URL,
            data={"text": text, "language": language},
            timeout=15,
        )
        result = response.json()

        issues = []
        for match in result.get("matches", []):
            issues.append({
                "message": match.get("message"),
                "offset": match.get("offset"),
                "length": match.get("length"),
                "suggestions": [r["value"] for r in match.get("replacements", [])[:3]],
            })

        return jsonify({
            "issues": issues,
            "total_issues": len(issues),
        })

    except requests.exceptions.RequestException as e:
        return jsonify({"error": f"Service unavailable: {str(e)}"}), 503


@app.route("/check-plagiarism", methods=["POST"])
def check_plagiarism():
    data = request.get_json()
    text = data.get("text", "")
    corpus = data.get("corpus", [])

    if not text.strip():
        return jsonify({"error": "Text is empty"}), 400
    if not corpus:
        return jsonify({"overall_score": 0, "matches": [], "note": "No corpus provided"})

    input_sentences = [s.strip() for s in text.split(".") if s.strip()]
    if not input_sentences:
        input_sentences = [text]

    input_embeddings = similarity_model.encode(input_sentences)
    corpus_embeddings = similarity_model.encode(corpus)

    matches = []
    total_similarity = 0

    for i, sent in enumerate(input_sentences):
        sims = cosine_similarity([input_embeddings[i]], corpus_embeddings)[0]
        best_match_idx = sims.argmax()
        best_score = float(sims[best_match_idx])

        if best_score > 0.6:
            matches.append({
                "sentence": sent,
                "matched_with": corpus[best_match_idx][:150],
                "similarity": round(best_score * 100, 1),
            })
        total_similarity += best_score

    overall_score = round((total_similarity / len(input_sentences)) * 100, 1)

    return jsonify({
        "overall_score": overall_score,
        "matches": matches,
        "sentences_checked": len(input_sentences),
    })


@app.route("/summarize", methods=["POST"])
def summarize():
    data = request.get_json()
    text = data.get("text", "")
    max_length = data.get("max_length", 150)

    if not text.strip():
        return jsonify({"error": "Text is empty"}), 400

    try:
        headers = {"Authorization": f"Bearer {HF_API_TOKEN}"}
        payload = {
            "inputs": text[:3000],
            "parameters": {"max_length": max_length, "min_length": 30},
        }
        response = requests.post(HF_SUMMARY_URL, headers=headers, json=payload, timeout=30)
        result = response.json()

        if isinstance(result, list) and len(result) > 0:
            summary_text = result[0].get("summary_text", "")
        else:
            return jsonify({"error": "Summarization failed", "details": result}), 503

        flashcards = generate_flashcards(summary_text)

        return jsonify({
            "summary": summary_text,
            "flashcards": flashcards,
        })

    except requests.exceptions.RequestException as e:
        return jsonify({"error": f"Service unavailable: {str(e)}"}), 503


def generate_flashcards(summary_text):
    sentences = [s.strip() for s in summary_text.split(".") if s.strip()]
    flashcards = []

    for sent in sentences[:5]:
        words = sent.split()
        if len(words) < 4:
            continue
        subject = " ".join(words[:3])
        flashcards.append({
            "question": f"Explain: {subject}...?",
            "answer": sent + ".",
        })

    return flashcards


def fetch_image_for_slide(query):
    """Pixabay se topic-related image dhoondta hai aur temporarily save karta hai"""
    if not PIXABAY_API_KEY:
        return None
    try:
        response = requests.get(
            "https://pixabay.com/api/",
            params={
                "key": PIXABAY_API_KEY,
                "q": query,
                "per_page": 3,
                "image_type": "photo",
                "safesearch": "true",
            },
            timeout=10,
        )
        result = response.json()

        if result.get("hits"):
            image_url = result["hits"][0]["webformatURL"]
            image_data = requests.get(image_url, timeout=10).content

            os.makedirs("temp_images", exist_ok=True)
            safe_name = "".join(c if c.isalnum() else "_" for c in query)[:30]
            image_path = f"temp_images/{safe_name}.jpg"

            with open(image_path, "wb") as f:
                f.write(image_data)

            return image_path
    except Exception as e:
        print(f"Image fetch error: {e}")

    return None


@app.route("/generate-slides", methods=["POST"])
def generate_slides():
    data = request.get_json()
    topic = data.get("topic", "")
    num_slides = data.get("num_slides", 5)
    theme = data.get("theme", "default")

    if not topic.strip():
        return jsonify({"error": "Topic is empty"}), 400

    prompt = f"""Generate content for a {num_slides}-slide presentation on the topic: "{topic}"

Return ONLY valid JSON in this exact format, no other text:
{{
  "title": "Presentation Title Here",
  "slides": [
    {{"heading": "Slide Heading", "bullets": ["point 1", "point 2", "point 3"]}}
  ]
}}

Make {num_slides} content slides. Keep bullets concise (under 12 words each)."""

    try:
        response = gemini_model.generate_content(prompt)
        text_response = response.text
        text_response = text_response.replace("```json", "").replace("```", "").strip()
        slide_data = json.loads(text_response)

    except Exception as e:
        print(f"GEMINI ERROR: {e}")
        slide_data = {
            "title": topic,
            "slides": [
                {"heading": "Introduction", "bullets": ["Overview of " + topic, "Why it matters", "What we'll cover"]},
                {"heading": "Key Concepts", "bullets": ["Main point one", "Main point two", "Main point three"]},
                {"heading": "Applications", "bullets": ["Real-world use case 1", "Real-world use case 2"]},
                {"heading": "Conclusion", "bullets": ["Summary", "Key takeaways", "Questions?"]},
            ]
        }

    output_dir = "generated_slides"
    os.makedirs(output_dir, exist_ok=True)
    safe_filename = "".join(c if c.isalnum() else "_" for c in topic)[:50]
    file_path = os.path.join(output_dir, f"{safe_filename}.pptx")

    create_pptx_file(slide_data, file_path, theme)

    return jsonify({
        "file_path": file_path,
        "slide_content": slide_data,
    })


def set_background_color(slide, color):
    """Slide ka background color set karta hai"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_pptx_file(slide_data, output_path, theme="default"):
    prs = Presentation()

    # ============================================
    # Themes - naye theme add karne ke liye bas
    # yahan ek nayi entry add karein (title, text, accent, bg colors ke sath)
    # ============================================
    themes = {
        "default": {
            "title": RGBColor(30, 30, 30),
            "text": RGBColor(60, 60, 60),
            "accent": RGBColor(220, 38, 38),
            "bg": RGBColor(255, 255, 255),
        },
        "dark": {
            "title": RGBColor(255, 255, 255),
            "text": RGBColor(220, 220, 220),
            "accent": RGBColor(96, 165, 250),
            "bg": RGBColor(25, 25, 35),
        },
        "ocean": {
            "title": RGBColor(15, 76, 117),
            "text": RGBColor(50, 50, 50),
            "accent": RGBColor(52, 152, 219),
            "bg": RGBColor(240, 248, 255),
        },
        "forest": {
            "title": RGBColor(27, 67, 50),
            "text": RGBColor(60, 60, 60),
            "accent": RGBColor(64, 145, 108),
            "bg": RGBColor(240, 250, 240),
        },
        "sunset": {
            "title": RGBColor(120, 40, 31),
            "text": RGBColor(70, 40, 30),
            "accent": RGBColor(230, 126, 34),
            "bg": RGBColor(255, 245, 235),
        },
    }
    colors = themes.get(theme, themes["default"])

    # ===== Title Slide (topic ki image ke saath) =====
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_background_color(slide, colors["bg"])

    slide.shapes.title.text = slide_data.get("title", "Presentation")
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = colors["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)

    title_image = fetch_image_for_slide(slide_data.get("title", ""))
    if title_image:
        slide.shapes.add_picture(title_image, Inches(5.5), Inches(1.5), width=Inches(4))

    # ===== Content Slides (har ek ki apni related image) =====
    bullet_layout = prs.slide_layouts[1]
    for slide_info in slide_data.get("slides", []):
        slide = prs.slides.add_slide(bullet_layout)
        set_background_color(slide, colors["bg"])

        slide.shapes.title.text = slide_info.get("heading", "")
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = colors["title"]

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        bullets = slide_info.get("bullets", [])
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = colors["text"]

        slide_image = fetch_image_for_slide(slide_info.get("heading", ""))
        if slide_image:
            slide.shapes.add_picture(slide_image, Inches(6), Inches(1.8), width=Inches(3.3))

    prs.save(output_path)


# ============================================
# File Serving Endpoint (slides download ke liye)
# ============================================
@app.route("/generated_slides/<filename>")
def serve_slide(filename):
    return send_from_directory("generated_slides", filename)


# ============================================
# Health Check Endpoint
# ============================================
@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "AI microservice running"})


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)